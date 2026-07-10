from __future__ import annotations

import argparse
import csv
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_DATASET_ROOT = PROJECT_ROOT / "dataset_new" / "Military_TrueFalse_Selected44_hard_projection"
DEFAULT_FIGURE_DIR = PROJECT_ROOT / "artifacts" / "figures" / "military_selected44_2026-07-06"
DEFAULT_OUTPUT_DIR = PROJECT_ROOT / "presentation_outputs"
GATE_RE = re.compile(r"^(?P<base>.+?)_gate_(?P<gate>\d+)\.png$")


@dataclass(frozen=True)
class GateStackSample:
    sample_id: str
    true_paths: dict[int, Path]
    false_paths: dict[int, Path]
    score: float


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build an editable PPTX for selected-44 military gated experiments.")
    parser.add_argument("--dataset-root", type=Path, default=DEFAULT_DATASET_ROOT)
    parser.add_argument("--figure-dir", type=Path, default=DEFAULT_FIGURE_DIR)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    parser.add_argument("--pptx-name", default="military_selected44_gated_report_2026-07-06.pptx")
    parser.add_argument("--notes-name", default="military_selected44_gated_report_speaker_notes_2026-07-06.md")
    return parser.parse_args()


def group_images(class_dir: Path) -> dict[str, dict[int, Path]]:
    grouped: dict[str, dict[int, Path]] = {}
    for path in sorted(class_dir.glob("*_gate_*.png")):
        match = GATE_RE.match(path.name)
        if not match:
            continue
        grouped.setdefault(match.group("base"), {})[int(match.group("gate"))] = path
    return grouped


def load_gray(path: Path) -> np.ndarray:
    return np.asarray(Image.open(path).convert("L"), dtype=np.float32)


def norm_image(arr: np.ndarray) -> np.ndarray:
    max_value = float(arr.max())
    if max_value <= 0.0:
        return np.zeros_like(arr, dtype=np.float32)
    return arr / max_value


def corr(a: np.ndarray, b: np.ndarray) -> float:
    af = a.reshape(-1)
    bf = b.reshape(-1)
    if float(af.std()) == 0.0 or float(bf.std()) == 0.0:
        return 0.0
    return float(np.corrcoef(af, bf)[0, 1])


def mean_pair_corr(paths: dict[int, Path]) -> float:
    arrays = [norm_image(load_gray(paths[idx])) for idx in sorted(paths)]
    values: list[float] = []
    for idx, left in enumerate(arrays):
        for right in arrays[idx + 1 :]:
            values.append(corr(left, right))
    return float(sum(values) / len(values)) if values else 0.0


def active_fraction(paths: dict[int, Path]) -> float:
    values = []
    for path in paths.values():
        arr = load_gray(path)
        values.append(float((arr > 5).mean()))
    return float(sum(values) / len(values)) if values else 0.0


def choose_sample(dataset_root: Path) -> GateStackSample:
    true_grouped = group_images(dataset_root / "true3d")
    false_grouped = group_images(dataset_root / "flat_false")
    candidates: list[GateStackSample] = []
    for sample_id, true_paths in sorted(true_grouped.items()):
        false_paths = false_grouped.get(sample_id)
        if not false_paths or len(true_paths) != 3 or len(false_paths) != 3:
            continue
        true_corr = mean_pair_corr(true_paths)
        false_corr = mean_pair_corr(false_paths)
        foreground = min(active_fraction(true_paths), active_fraction(false_paths))
        score = (false_corr - true_corr) + foreground
        candidates.append(GateStackSample(sample_id, true_paths, false_paths, score))
    if not candidates:
        raise RuntimeError(f"No paired true3d/flat_false 3-gate samples found under {dataset_root}")
    return max(candidates, key=lambda item: item.score)


def make_gate_stack_figure(dataset_root: Path, output_dir: Path) -> Path:
    sample = choose_sample(dataset_root)
    out_path = output_dir / "true3d_vs_hard_projection_gate_stack.png"
    output_dir.mkdir(parents=True, exist_ok=True)

    fig, axes = plt.subplots(2, 3, figsize=(8.4, 4.4))
    rows = [("True 3D target", sample.true_paths), ("Hard projection false target", sample.false_paths)]
    for row_idx, (row_label, paths) in enumerate(rows):
        for gate in range(3):
            ax = axes[row_idx, gate]
            arr = load_gray(paths[gate])
            ax.imshow(arr, cmap="gray", vmin=0, vmax=max(20.0, float(arr.max())))
            ax.set_xticks([])
            ax.set_yticks([])
            if row_idx == 0:
                ax.set_title(f"gate_{gate}", fontsize=11)
            if gate == 0:
                ax.set_ylabel(row_label, fontsize=10)
    fig.suptitle("True 3D vs Hard Projection Gate Stack", fontsize=14)
    fig.tight_layout(rect=(0, 0, 1, 0.95))
    fig.savefig(out_path, dpi=240)
    plt.close(fig)
    return out_path


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=(39, 49, 58), align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.45), title, 25, True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.75), Inches(11.5), Inches(0.35), subtitle, 12, False, (93, 103, 113))


def add_bullets(slide, left, top, width, height, bullets: list[str], font_size=16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(39, 49, 58)
        p.space_after = Pt(8)


def add_metric_cards(slide, metrics: list[tuple[str, str, str]], top=1.45) -> None:
    left = Inches(0.65)
    width = Inches(3.85)
    gap = Inches(0.25)
    for idx, (value, label, note) in enumerate(metrics):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + idx * (width + gap),
            Inches(top),
            width,
            Inches(1.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(242, 246, 250)
        shape.line.color.rgb = RGBColor(196, 205, 214)
        add_textbox(slide, left + idx * (width + gap) + Inches(0.15), Inches(top + 0.12), width - Inches(0.3), Inches(0.3), value, 21, True, (31, 82, 130), PP_ALIGN.CENTER)
        add_textbox(slide, left + idx * (width + gap) + Inches(0.15), Inches(top + 0.47), width - Inches(0.3), Inches(0.22), label, 11, True, (39, 49, 58), PP_ALIGN.CENTER)
        add_textbox(slide, left + idx * (width + gap) + Inches(0.15), Inches(top + 0.72), width - Inches(0.3), Inches(0.2), note, 9, False, (93, 103, 113), PP_ALIGN.CENTER)


def add_image(slide, path: Path, left, top, width=None, height=None) -> None:
    if not path.exists():
        raise FileNotFoundError(path)
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def create_pptx(figure_paths: dict[str, Path], output_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.8), "基于激光距离选通序列的军事三维目标识别与二维假目标判别", 27, True)
    add_textbox(slide, Inches(0.72), Inches(2.35), Inches(10.8), Inches(0.35), "精选 44 个军事 3D 模型 | Blender gate 仿真 | attention_residual 网络 | hard projection 假目标", 15, False, (79, 89, 99))
    add_metric_cards(
        slide,
        [
            ("44", "人工筛选模型", "tank / fighter / helicopter"),
            ("3 gates", "距离选通序列", "gate_0, gate_1, gate_2"),
            ("100%", "hard projection full stack", "三随机种子 mean best val acc"),
        ],
        top=3.25,
    )

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "数据治理：只使用人工 keep=1 的军事模型", "来源：thumbnail_review.csv；不直接使用噪声较大的候选池")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(5.1),
        Inches(2.6),
        [
            "01_Main_Battle_Tank：12 个",
            "02_Fighter_Jet：20 个",
            "03_Attack_Helicopter：12 个",
            "合计：44 个高质量精选模型",
        ],
        17,
    )
    add_bullets(
        slide,
        Inches(6.2),
        Inches(1.35),
        Inches(6.2),
        Inches(3.8),
        [
            "Blender 生成 true3d 三门 gate stack",
            "flat-echo 生成整目标轮廓强度衰减假目标",
            "hard projection 从真实 3D gate stack 最大投影生成更难假目标",
            "所有训练使用 split_group_by_sample_id，避免同源模型泄漏",
        ],
        16,
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "方法：共享 CNN 编码每个 gate，再做 gate-level 融合")
    add_bullets(
        slide,
        Inches(0.7),
        Inches(1.25),
        Inches(5.8),
        Inches(4.5),
        [
            "输入：同一目标的 [gate_0, gate_1, gate_2]",
            "SliceEncoder：共享 CNN，逐 gate 提取特征",
            "Fusion：attention_residual 融合多 gate 特征",
            "输出：军事三分类或 true3d / flat_false 二分类",
            "attention 解释为 gate-level discriminative contribution",
        ],
        15,
    )
    add_image(slide, figure_paths["gate_stack"], Inches(6.35), Inches(1.18), width=Inches(6.25))

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "军事三分类：迁移学习提高小样本稳定性", "44 个样本条件下，稳定性比单次最高精度更重要")
    add_image(slide, figure_paths["transfer"], Inches(0.75), Inches(1.15), width=Inches(7.05))
    add_bullets(
        slide,
        Inches(8.15),
        Inches(1.35),
        Inches(4.45),
        Inches(3.8),
        [
            "Transfer frozen / finetune：mean = 0.7500，std = 0",
            "Scratch：mean = 0.7083，std = 0.1443",
            "结论：迁移学习当前不是显著提高上限，而是降低小样本随机波动",
        ],
        15,
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "核心证据：hard projection 下单 gate 接近随机，完整 gate stack 稳定可分")
    add_image(slide, figure_paths["hard_ablation"], Inches(0.75), Inches(1.1), width=Inches(7.55))
    add_bullets(
        slide,
        Inches(8.55),
        Inches(1.25),
        Inches(4.1),
        Inches(4.2),
        [
            "hard projection 假目标使用真实 3D gate stack 最大投影作为二维轮廓",
            "单 gate：53.7% / 63.0% / 53.7%",
            "完整三门：100.0%",
            "结论：判别信息主要来自跨 gate 序列响应，而不是普通单帧外观",
        ],
        15,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "物理诊断：平面假目标跨 gate 高相似，真实三维目标跨 gate 低相似")
    add_image(slide, figure_paths["diagnostics"], Inches(0.65), Inches(1.02), width=Inches(7.85))
    add_bullets(
        slide,
        Inches(8.75),
        Inches(1.2),
        Inches(3.9),
        Inches(4.3),
        [
            "flat_false：同一整目标轮廓在多个 gate 中强度缩放",
            "true3d：不同 gate 响应不同深度结构",
            "hard projection flat_false：corr = 0.9768，IoU = 0.9301",
            "hard projection true3d：corr = 0.3246，IoU = 0.3065",
        ],
        14,
    )

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "鲁棒性：对噪声较稳，对 gate 缺失更敏感")
    add_image(slide, figure_paths["robustness"], Inches(0.72), Inches(1.12), width=Inches(7.6))
    add_bullets(
        slide,
        Inches(8.55),
        Inches(1.3),
        Inches(4.1),
        Inches(3.9),
        [
            "Clean full stack：100.0%",
            "Noise + background + Poisson：98.1%",
            "Random gate dropout：90.7%",
            "结论：成像噪声不是主要瓶颈；gate stack 完整性需要保证",
        ],
        15,
    )

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "阶段结论与下一步")
    add_bullets(
        slide,
        Inches(0.72),
        Inches(1.2),
        Inches(5.95),
        Inches(4.8),
        [
            "已完成精选军事 44 模型子集、真/假目标 gate stack、训练就绪检查与多组训练验证",
            "hard projection 结果直接支撑：激光选通序列提供了单帧图像不具备的判别信息",
            "当前电子网络可作为后续多模光神经网络高速识别的可验证基线",
        ],
        15,
    )
    add_bullets(
        slide,
        Inches(7.05),
        Inches(1.2),
        Inches(5.55),
        Inches(4.8),
        [
            "继续扩大高质量军事模型数量",
            "补充 num gates = 1 / 3 / 5、gate width、pulse width 消融",
            "加入更真实材质反射率、姿态变化和复杂背景",
            "设计光学编码/多模光纤散斑前端接口",
        ],
        15,
    )

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)


def write_notes(path: Path) -> None:
    text = """# 军事 44 模型 gate stack 汇报讲稿

## 1. 研究问题
本页强调：普通二维图像只能看到目标投影，而激光距离选通可以得到同一目标在多个深度门下的响应序列。本文要验证的是，这个序列是否能帮助区分真实三维目标和平面二维假目标。

## 2. 数据治理与仿真流程
说明军事数据没有直接全量使用，而是通过 thumbnail_review.csv 人工筛选 keep=1 的 44 个模型。这个步骤体现数据治理，避免标签噪声直接污染训练。

## 3. 网络结构
说明输入是三张 gate 图像，不是 RGB。每张 gate 用共享 CNN 编码，再用 attention_residual 融合。attention 只解释为门控切片级判别贡献。

## 4. 军事三分类迁移学习
重点讲稳定性：迁移学习 mean=0.75 且 std=0，从零训练 mean=0.7083 且 std=0.1443。当前结论是迁移让小样本训练更稳定，不夸大为显著提高上限。

## 5. hard projection 核心结果
这是主结果。hard projection 用真实 3D gate stack 的最大投影生成平面诱饵，压低单帧图像差异。单 gate 接近随机，完整三门达到 100%，直接证明选通序列的必要性。

## 6. 物理诊断
平面假目标跨 gate 高相关、高 IoU；真实三维目标跨 gate 低相关、低 IoU。这个诊断把网络结果和距离选通物理机制对应起来。

## 7. 鲁棒性
噪声、背景、Poisson 退化下仍较稳，随机丢门下降。结论是 gate 序列完整性比中等强度噪声更关键。

## 8. 总结与下一步
收束到三点：数据筛选与仿真链路已跑通；hard projection 支撑选通序列价值；当前电子网络是后续光神经网络融合的基线。
"""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def inspect_pptx(path: Path) -> tuple[int, list[str]]:
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        snippets = []
        for name in slide_names[:3]:
            xml = zf.read(name).decode("utf-8", errors="ignore")
            snippets.append(name if "基于激光距离选通序列" in xml or "数据治理" in xml or "方法" in xml else f"{name}: text-present")
    return len(slide_names), snippets


def main() -> int:
    args = parse_args()
    figure_dir = args.figure_dir
    gate_stack = make_gate_stack_figure(args.dataset_root, figure_dir)
    figure_paths = {
        "gate_stack": gate_stack,
        "transfer": figure_dir / "military_3class_transfer_vs_scratch.png",
        "hard_ablation": figure_dir / "hard_projection_full_stack_vs_single_gate.png",
        "diagnostics": figure_dir / "gate_stack_physical_diagnostics.png",
        "robustness": figure_dir / "per_gate_norm_robustness.png",
    }
    output_pptx = args.output_dir / args.pptx_name
    output_notes = args.output_dir / args.notes_name
    create_pptx(figure_paths, output_pptx)
    write_notes(output_notes)
    slide_count, snippets = inspect_pptx(output_pptx)
    print(f"pptx={output_pptx}")
    print(f"notes={output_notes}")
    print(f"gate_stack_figure={gate_stack}")
    print(f"slide_count={slide_count}")
    for snippet in snippets:
        print(snippet)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
